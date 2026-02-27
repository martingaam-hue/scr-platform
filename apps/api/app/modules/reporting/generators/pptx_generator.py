"""PPTX report generator using python-pptx."""

import io

from pptx import Presentation
from pptx.util import Inches, Pt

from app.modules.reporting.generators.base import BaseReportGenerator


class PPTXGenerator(BaseReportGenerator):
    """Generate PowerPoint presentations with title and data slides."""

    CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def generate(self, data: dict, sections: list[dict]) -> tuple[bytes, str]:
        prs = Presentation()
        self._create_title_slide(prs, data)

        for section in sections:
            self._create_section_slide(prs, section, data)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue(), self.CONTENT_TYPE

    def _create_title_slide(self, prs: Presentation, data: dict) -> None:
        layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = data.get("title", "Report")
        if slide.placeholders[1]:
            slide.placeholders[1].text = (
                f"{self.org_name}\n{self.generated_at}"
            )

    def _create_section_slide(
        self, prs: Presentation, section: dict, data: dict
    ) -> None:
        name = section.get("name", "Section")
        section_data = data.get(name, {})

        layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = name.replace("_", " ").title()

        if isinstance(section_data, list) and section_data:
            # Create table
            headers = list(section_data[0].keys())
            rows_count = min(len(section_data), 15)  # Limit rows per slide
            cols_count = min(len(headers), 6)

            table = slide.shapes.add_table(
                rows_count + 1,
                cols_count,
                Inches(0.5),
                Inches(1.5),
                Inches(9),
                Inches(0.4 * (rows_count + 1)),
            ).table

            for col_idx in range(cols_count):
                cell = table.cell(0, col_idx)
                cell.text = headers[col_idx].replace("_", " ").title()
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(10)

            for row_idx, row_data in enumerate(section_data[:rows_count], 1):
                for col_idx in range(cols_count):
                    cell = table.cell(row_idx, col_idx)
                    val = row_data.get(headers[col_idx], "")
                    cell.text = str(val) if val is not None else ""
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(9)

        elif isinstance(section_data, dict):
            # Key-value pairs as bullet points
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for key, val in section_data.items():
                p = tf.add_paragraph()
                p.text = f"{key.replace('_', ' ').title()}: {val}"
                p.font.size = Pt(12)

        else:
            body = slide.placeholders[1]
            body.text = str(section_data) if section_data else f"No data for {name}"
